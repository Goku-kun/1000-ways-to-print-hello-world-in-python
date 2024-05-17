#pip install python-pptx
from pptx import Presentation

# Load the PowerPoint presentation
presentation = Presentation('using-powerpoint.pptx')

# Get the first slide
first_slide = presentation.slides[0]

# Get the title of the first slide
title = first_slide.shapes.title.text

print(title)